import os
import sys
import re
from pptx import Presentation

def parse_markdown(filepath):
    with open(filepath, 'r', encoding='utf-8') as f:
        content = f.read()

    # Split by '---' to get slides.
    raw_slides = re.split(r'\n---\n', content)
    slides_data = []

    for raw in raw_slides:
        raw = raw.strip()
        # Skip frontmatter block
        if not raw or raw.startswith('marp:') or raw.startswith('theme:'):
            continue
            
        slide = {
            'layout': 1, # Default to "Title and Content"
            'title': '',
            'content': [],
            'notes': ''
        }
        
        # Identify slide class to pick the right PowerPoint Slide Layout
        class_match = re.search(r'<!--\s*_class:\s*(.*?)\s*-->', raw)
        if class_match:
            c = class_match.group(1).strip()
            if c == 'slide-cover':
                slide['layout'] = 0 # Title Slide
            elif c == 'slide-divider':
                slide['layout'] = 2 # Section Header
            elif c == 'slide-simple':
                slide['layout'] = 1 # Title and Content
                
        # Extract speaker notes
        notes_match = re.search(r'<!--\s*_speaker_notes:\s*(.*?)\s*-->', raw, re.DOTALL)
        if notes_match:
            slide['notes'] = notes_match.group(1).strip()
            
        # Parse title and content
        lines = raw.split('\n')
        for line in lines:
            line = line.strip()
            if line.startswith('# '):
                if not slide['title']:
                    slide['title'] = line[2:].strip()
            elif line.startswith('## ') or line.startswith('### '):
                # Use subheaders as title if not yet set, otherwise add to content
                if not slide['title']:
                    slide['title'] = line.lstrip('# ').strip()
                else:
                    slide['content'].append(line.lstrip('# ').strip())
            elif line.startswith('- ') or line.startswith('* '):
                slide['content'].append(line[2:].strip())
            elif line and not line.startswith('<') and not line.startswith('!['):
                # Regular text ignoring HTML comments and images (images would need complex placement)
                slide['content'].append(line.strip())
                
        slides_data.append(slide)
        
    return slides_data

def create_pptx(slides_data, template_path, output_path):
    prs = Presentation(template_path)
    
    for data in slides_data:
        try:
            layout = prs.slide_layouts[data['layout']]
        except IndexError:
            # Fallback to layout 1 if index doesn't exist
            layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
            
        slide = prs.slides.add_slide(layout)
        
        # Populate Title
        if slide.shapes.title and data['title']:
            slide.shapes.title.text = data['title']
            
        # Populate Content
        # Usually, placeholders[1] is the main content body.
        body_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                body_placeholder = shape
                break
        
        # Fallback if idx 1 is not found
        if not body_placeholder and len(slide.placeholders) > 1:
            body_placeholder = slide.placeholders[1]
            
        if body_placeholder and data['content']:
            tf = body_placeholder.text_frame
            tf.text = data['content'][0]
            for idx in range(1, len(data['content'])):
                p = tf.add_paragraph()
                p.text = data['content'][idx]
                
        # Populate Speaker Notes
        if data['notes']:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = data['notes']
            
    prs.save(output_path)
    print(f"Generated Editable PPTX saved to {output_path}")

if __name__ == '__main__':
    if len(sys.argv) < 3:
        print("Usage: python md_to_native_pptx.py <input.md> <output.pptx> [template.pptx]")
        sys.exit(1)
        
    input_md = sys.argv[1]
    output_pptx = sys.argv[2]
    
    # Default template path relative to script location
    script_dir = os.path.dirname(os.path.abspath(__file__))
    default_template = os.path.join(script_dir, '..', '..', 'assets', 'mindx-b2b-master.pptx')
    
    template_pptx = sys.argv[3] if len(sys.argv) > 3 else default_template
    
    if not os.path.exists(template_pptx):
        print(f"Error: Master template not found at {template_pptx}")
        sys.exit(1)
        
    print(f"Using template: {template_pptx}")
    data = parse_markdown(input_md)
    print(f"Parsed {len(data)} slides.")
    create_pptx(data, template_pptx, output_pptx)
