from pptx import Presentation

# Create a presentation based on default template
prs = Presentation()

# The default template has 11 slide layouts, we don't necessarily need to add slides.
# Just saving the template is enough. But let's add one slide to see if it works.
slide_layout = prs.slide_layouts[0] # Title slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "MindX B2B Master Template"
subtitle.text = "Replace this file with the official designer's template."

# Save to the assets folder
prs.save('.agents/skills/s-chuan-hoa-tai-lieu/assets/mindx-b2b-master.pptx')
print("Master template saved successfully.")
